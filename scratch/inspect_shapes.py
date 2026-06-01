from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import sys
import io

sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

def inspect_shapes(filename, slide_num):
    prs = Presentation(filename)
    slide = prs.slides[slide_num - 1]
    print(f"\n--- Slide {slide_num} ---")
    for shape in slide.shapes:
        print(f"Shape: {shape.name}, Type: {shape.shape_type}")
        if shape.has_text_frame:
            print(f"  Text: {shape.text}")
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            print("  Group Shape Contents:")
            for s in shape.shapes:
                print(f"    Sub-Shape: {s.name}, Type: {s.shape_type}")
                if s.has_text_frame:
                    print(f"      Text: {s.text}")

inspect_shapes('데이터 아키텍처 정의.pptx', 359)
