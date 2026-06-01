from pptx import Presentation
import sys

# Set encoding for output to UTF-8
import io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

def extract_slides(filename, start, end):
    prs = Presentation(filename)
    for i, slide in enumerate(prs.slides):
        if i + 1 < start:
            continue
        if i + 1 > end:
            break
        print(f"\n--- Slide {i + 1} ---")
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                print(shape.text)
            if shape.has_table:
                for row in shape.table.rows:
                    print(" | ".join(cell.text_frame.text.replace("\n", " ") for cell in row.cells))

extract_slides('데이터 아키텍처 정의.pptx', 350, 400)
