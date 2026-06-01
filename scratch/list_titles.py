from pptx import Presentation
import sys
import io

sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

def list_titles(filename, start, end):
    prs = Presentation(filename)
    for i, slide in enumerate(prs.slides):
        if i + 1 < start:
            continue
        if i + 1 > end:
            break
        title = "No Title"
        if slide.shapes.title:
            title = slide.shapes.title.text
        print(f"Slide {i+1}: {title}")

list_titles('데이터 아키텍처 정의.pptx', 260, 360)
