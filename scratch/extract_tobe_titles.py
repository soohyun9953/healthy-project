from pptx import Presentation
import sys, io

sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')
prs = Presentation('To-Be 업무프로세스_v0.561(F).pptx')

print(f"총 슬라이드 수: {len(prs.slides)}")
print()
for i, slide in enumerate(prs.slides):
    snum = i + 1
    title_text = ""
    for shape in slide.shapes:
        if shape.has_text_frame and shape.shape_type in (13, 14, 1):
            t = shape.text_frame.text.strip()
            if t and len(t) < 120:
                title_text = t[:100]
                break
    if not title_text:
        for shape in slide.shapes:
            if shape.has_text_frame:
                t = shape.text_frame.text.strip()
                if t and len(t) < 120:
                    title_text = t[:100]
                    break
    print(f"Slide {snum:3d}: {title_text}")
