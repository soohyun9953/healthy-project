from pptx import Presentation
from pptx.util import Pt
import sys

prs = Presentation('To-Be 업무프로세스_v0.561(F).pptx')

# 슬라이드 노트 및 도형에서 모든 텍스트 수집
# 패턴을 파악하기 위해 처음 5개, 특별 슬라이드 몇 개 상세 추출
special_slides = [1, 2, 3, 95, 257, 270, 350, 400, 500, 600, 626]

for snum in special_slides:
    slide = prs.slides[snum - 1]
    print(f"\n{'='*80}")
    print(f"[Slide {snum}]")
    print('='*80)
    for shape in slide.shapes:
        if shape.has_text_frame:
            txt = shape.text_frame.text.strip()
            if txt:
                print(f"  [SHAPE:{shape.shape_type}] {txt[:300]}")
        if hasattr(shape, 'table'):
            tbl = shape.table
            for row in tbl.rows:
                row_texts = [c.text_frame.text.strip() for c in row.cells if c.text_frame.text.strip()]
                if row_texts:
                    print(f"  [TABLE] {' | '.join(row_texts[:6])}")
