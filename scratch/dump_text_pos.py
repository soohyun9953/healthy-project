from pptx import Presentation
import sys
import io

sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

def dump_all_text_with_pos(filename, slide_num):
    prs = Presentation(filename)
    slide = prs.slides[slide_num - 1]
    print(f"\n--- Slide {slide_num} ---")
    for shape in slide.shapes:
        if hasattr(shape, "text"):
            print(f"Shape: {shape.name}, Pos: ({shape.left}, {shape.top}), Text: {shape.text}")
        if shape.has_table:
            print(f"Table Found: {shape.name}")
            for r_idx, row in enumerate(shape.table.rows):
                for c_idx, cell in enumerate(row.cells):
                    if cell.text_frame.text.strip():
                        print(f"  Cell({r_idx}, {c_idx}): {cell.text_frame.text}")

dump_all_text_with_pos('데이터 아키텍처 정의.pptx', 359)
