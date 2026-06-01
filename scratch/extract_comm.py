from pptx import Presentation
import sys
import io

sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

def extract_range(filename, start, end):
    prs = Presentation(filename)
    for i, slide in enumerate(prs.slides):
        if i + 1 < start:
            continue
        if i + 1 > end:
            break
        print(f"\n--- Slide {i + 1} ---")
        if slide.shapes.title:
            print(f"Title: {slide.shapes.title.text}")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape != slide.shapes.title:
                print(f"Text: {shape.text}")
            if shape.has_table:
                print("Table Found:")
                for row in shape.table.rows:
                    print(" | ".join(cell.text_frame.text.replace("\n", " ") for cell in row.cells))

extract_range('데이터 아키텍처 정의.pptx', 43, 72)
