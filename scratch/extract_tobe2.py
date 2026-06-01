from pptx import Presentation
import io, sys

prs = Presentation('To-Be 업무프로세스_v0.561(F).pptx')
print(f"총 슬라이드 수: {len(prs.slides)}", flush=True)

for i, slide in enumerate(prs.slides):
    snum = i + 1
    texts = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            t = shape.text_frame.text.strip()
            if t:
                texts.append(t[:80])
    title = texts[0] if texts else "(내용없음)"
    print(f"S{snum:4d}: {title[:100]}", flush=True)
