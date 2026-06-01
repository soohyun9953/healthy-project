from pptx import Presentation
import sys
import io

sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

def extract_columns(filename, start, end):
    prs = Presentation(filename)
    for i, slide in enumerate(prs.slides):
        if i + 1 < start:
            continue
        if i + 1 > end:
            break
        print(f"\n--- Slide {i + 1} ---")
        if slide.shapes.title:
            print(f"Title: {slide.shapes.title.text}")
        for shape in slide.shapes:
            if shape.has_table:
                for r_idx, row in enumerate(shape.table.rows):
                    row_data = []
                    for c_idx, cell in enumerate(row.cells):
                        cell_text = "".join(p.text for p in cell.text_frame.paragraphs)
                        row_data.append(cell_text.strip())
                    print(" | ".join(row_data))

extract_columns('데이터 아키텍처 정의.pptx', 96, 110)
