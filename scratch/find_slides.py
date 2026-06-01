from pptx import Presentation
import sys
import io

sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

def find_slides(filename, keywords):
    prs = Presentation(filename)
    for i, slide in enumerate(prs.slides):
        title = ""
        if slide.shapes.title:
            title = slide.shapes.title.text
        
        content = ""
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                content += shape.text + " "
        
        if any(kw in title for kw in keywords) or any(kw in content for kw in keywords):
            print(f"Slide {i+1}: {title}")
            # Print first 100 chars of content
            print(f"  Content: {content[:100]}...")

find_slides('데이터 아키텍처 정의.pptx', ['커뮤니케이션', '협업', '논리 모델', '주제영역'])
