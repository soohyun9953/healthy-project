from pptx import Presentation
import re, io

prs = Presentation('To-Be 업무프로세스_v0.561(F).pptx')

out = []
out.append(f"총 슬라이드 수: {len(prs.slides)}")

for i, slide in enumerate(prs.slides):
    snum = i + 1
    all_texts = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            t = shape.text_frame.text.strip()
            if t:
                all_texts.append(t)
        if hasattr(shape, 'table'):
            tbl = shape.table
            for row in tbl.rows:
                row_txts = [c.text_frame.text.strip() for c in row.cells if c.text_frame.text.strip()]
                if row_txts:
                    all_texts.append('[TABLE] ' + ' | '.join(row_txts[:8]))

    # 제목 추출 (슬라이드 제목 플레이스홀더 우선)
    title = ""
    for shape in slide.shapes:
        if shape.shape_type == 14 and shape.has_text_frame:  # PLACEHOLDER
            t = shape.text_frame.text.strip()
            if t:
                title = t
                break
    if not title and all_texts:
        title = all_texts[0]

    combined = '\n    '.join(all_texts)
    safe = combined.encode('cp949', errors='replace').decode('cp949')
    out.append(f"\n[S{snum:4d}] {title[:80]}")
    out.append(f"    {safe[:500]}")

with open('tobe_full_extract.txt', 'w', encoding='utf-8') as f:
    f.write('\n'.join(out))
print("Done")
