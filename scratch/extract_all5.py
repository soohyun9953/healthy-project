from pptx import Presentation
import sys, io

sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')
prs = Presentation('데이터 아키텍처 정의.pptx')

ranges = {
    '바이오개체': (227, 262),
    '자원재고': (263, 300),
    '교육지식': (301, 321),
    'AI분석': (322, 343),
    '시스템운영': (344, 358),
}

for name, (start, end) in ranges.items():
    print(f"\n{'='*60}")
    print(f"[{name}] Slide {start}-{end}")
    print('='*60)
    for i, slide in enumerate(prs.slides):
        snum = i + 1
        if snum < start or snum > end:
            continue
        for shape in slide.shapes:
            if shape.has_table:
                for row in shape.table.rows:
                    cells = ["".join(p.text for p in cell.text_frame.paragraphs).strip()
                             for cell in row.cells]
                    print(f"S{snum}| " + " | ".join(cells))
