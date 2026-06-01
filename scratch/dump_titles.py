from pptx import Presentation
import sys
import io

sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

prs = Presentation('데이터 아키텍처 정의.pptx')
with open('slide_titles.txt', 'w', encoding='utf-8') as f:
    for i, slide in enumerate(prs.slides):
        title = "No Title"
        if slide.shapes.title:
            title = slide.shapes.title.text
        f.write(f"Slide {i+1}: {title}\n")
