from pptx import Presentation
import sys

prs = Presentation('정보시스템 응용 아키텍처 정의.pptx')
print(f"총 슬라이드 수: {len(prs.slides)}", flush=True)

for i, slide in enumerate(prs.slides):
    snum = i + 1
    texts = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            t = shape.text_frame.text.strip()
            if t:
                texts.append(t)
        if hasattr(shape, 'table'):
            tbl = shape.table
            for row in tbl.rows:
                row_txts = [c.text_frame.text.strip() for c in row.cells if c.text_frame.text.strip()]
                if row_txts:
                    texts.append('[TABLE] ' + ' | '.join(row_txts[:8]))

    title = ""
    for shape in slide.shapes:
        if shape.shape_type == 14 and shape.has_text_frame:
            t = shape.text_frame.text.strip()
            if t:
                title = t
                break
    if not title and texts:
        title = texts[0]

    combined = '\n    '.join(texts)
    safe = combined.encode('cp949', errors='replace').decode('cp949')
    print(f"\n[S{snum:4d}] {title[:100]}", flush=True)
    print(f"    {safe[:600]}", flush=True)
