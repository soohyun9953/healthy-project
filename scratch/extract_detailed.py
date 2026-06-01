from pptx import Presentation
import sys
import io

sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

def extract_table_detailed(filename, slide_num):
    prs = Presentation(filename)
    slide = prs.slides[slide_num - 1]
    print(f"\n--- Slide {slide_num} ---")
    for shape in slide.shapes:
        if shape.has_table:
            print("Table Found:")
            for r_idx, row in enumerate(shape.table.rows):
                row_data = []
                for c_idx, cell in enumerate(row.cells):
                    cell_text = ""
                    for paragraph in cell.text_frame.paragraphs:
                        for run in paragraph.runs:
                            cell_text += run.text
                    row_data.append(cell_text.strip())
                print(f"Row {r_idx}: " + " | ".join(row_data))

extract_table_detailed('데이터 아키텍처 정의.pptx', 43)
extract_table_detailed('데이터 아키텍처 정의.pptx', 44)
